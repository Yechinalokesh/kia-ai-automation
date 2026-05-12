import io, math, datetime
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Colors ─────────────────────────────────────────────────────────────────────
BLK  = RGBColor(0x00,0x00,0x00)
WHT  = RGBColor(0xFF,0xFF,0xFF)
NAV  = RGBColor(0x1F,0x38,0x64)
DRK  = RGBColor(0x26,0x26,0x26)
BLU  = RGBColor(0x20,0x60,0xA0)
RED  = RGBColor(0xC0,0x00,0x00)
GRN  = RGBColor(0x37,0x56,0x23)
YLW  = RGBColor(0xFF,0xFF,0xCC)
GRY1 = RGBColor(0xF2,0xF2,0xF2)
LBBL = RGBColor(0xDD,0xEA,0xF5)
BDR  = RGBColor(0xBF,0xBF,0xBF)
GRN2 = RGBColor(0xA9,0xD1,0x8E)

def I(v): return Inches(v)
def P(v): return Pt(v)

# ── Excel reader ───────────────────────────────────────────────────────────────
def read_excel(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    ws = wb.active

    def gv(c):
        v = ws[c].value
        if v is None: return None
        if isinstance(v, float) and math.isnan(v): return None
        if isinstance(v, str) and v.startswith('#'): return None
        return v

    def gn(c):
        v = gv(c)
        try: return int(float(v)) if v is not None else 0
        except (ValueError, TypeError): return 0

    def gf(c):
        v = gv(c)
        try: return float(v) if v is not None else 0.0
        except (ValueError, TypeError): return 0.0

    date_raw = gv('R3')
    if isinstance(date_raw, datetime.datetime):
        d2 = date_raw.day
        sfx = "th" if 11<=d2<=13 else {1:"st",2:"nd",3:"rd"}.get(d2%10,"th")
        date_str = f"\u300c{d2}{sfx} {date_raw.strftime('%b %Y')}\u300d"
        date_obj = date_raw
    else:
        date_str = str(date_raw) if date_raw else "N/A"
        date_obj = None

    def pct_str(v):
        try:
            f = float(v)
            return f"{f*100:.0f}%"
        except: return "0%"

    def ratio_pct(out_m, in_m):
        try:
            v = out_m/in_m*100 if in_m else 0
            return f"{v:.1f}"
        except: return "-"

    eff_c = gf('S12')

    k10t_in_m  = gn('BF43'); k10t_out_m  = gn('BF50')
    ak10m_in_m = gn('BF44'); ak10m_out_m = gn('BF51')
    ak12m_in_m = gn('BF45'); ak12m_out_m = gn('BF52')
    ak12c_in_m = gn('BF46'); ak12c_out_m = gn('BF53')
    gege_in_m  = gn('BF47'); gege_out_m  = gn('BF54')
    geng_in_m  = gn('BF48'); geng_out_m  = gn('BF55')

    ttl_in_m  = k10t_in_m+ak10m_in_m+ak12m_in_m+ak12c_in_m+gege_in_m+geng_in_m
    ttl_out_m = k10t_out_m+ak10m_out_m+ak12m_out_m+ak12c_out_m+gege_out_m+geng_out_m

    return {
        "date_str": date_str, "date_obj": date_obj,
        "uph": gn('O8'),
        "pt_c":gn('D12'),"pt_a":gn('E12'),"pt_b":gn('F12'),
        "apt_c":gn('N12'),"apt_a":gn('O12'),"apt_b":gn('P12'),
        "apt_month":gn('R12'),
        "eff_c": pct_str(eff_c), "eff_c_val": round(eff_c*100,1) if eff_c else 100.0,
        "eff_m": pct_str(gf('V12')),
        # Input C shift
        "k10t_in_c":gn('E18'),"ak10m_in_c":gn('E19'),
        "ak12m_in_c":gn('E20'),"gege_in_c":gn('E22'),"geng_in_c":gn('E23'),
        "ttl_in_c":gn('E24'),
        # Output C shift
        "k10t_out_c":gn('J18'),"ttl_out_c":gn('J24'),
        # Stock
        "k10t_stk":gn('Z18'),"ak10m_stk":gn('Z19'),"ak12m_stk":gn('Z20'),
        "gege_stk":gn('Z22'),"geng_stk":gn('Z23'),"ttl_stk":gn('Z24'),
        # Monthly totals
        "k10t_in_m":k10t_in_m,"k10t_out_m":k10t_out_m,
        "ak10m_in_m":ak10m_in_m,"ak10m_out_m":ak10m_out_m,
        "ak12m_in_m":ak12m_in_m,"ak12m_out_m":ak12m_out_m,
        "gege_in_m":gege_in_m,"gege_out_m":gege_out_m,
        "geng_in_m":geng_in_m,"geng_out_m":geng_out_m,
        "ttl_in_m":ttl_in_m,"ttl_out_m":ttl_out_m,
        # % calcs
        "k10t_pct":  ratio_pct(k10t_out_m, k10t_in_m),
        "ak12m_pct": ratio_pct(ak12m_out_m,ak12m_in_m),
        "gege_pct":  ratio_pct(gege_out_m, gege_in_m),
        "geng_pct":  ratio_pct(geng_out_m, geng_in_m),
        "ttl_pct":   ratio_pct(ttl_out_m,  ttl_in_m),
        # Feb plan/result from month cols
        "k10t_feb_p":gn('I18'),"k10t_feb_r":gn('N18'),
        "ak12m_feb_p":gn('I20'),"ak12m_feb_r":gn('N20'),
        "gege_feb_p":gn('I22'),"gege_feb_r":gn('N22'),
        "geng_feb_p":gn('I23'),"geng_feb_r":gn('N23'),
        "ttl_feb_p":gn('I24'),"ttl_feb_r":gn('N24'),
        # Down time
        "dt_equip":gn('C32'),"dt_nopart":gn('F32'),
        "dt_ngpart":gn('I32'),"dt_wdlay":gn('K32'),
        "dt_mcset":gn('M32'),"dt_plan":gn('O32'),
        "dt_other":gn('R32'),"dt_ttl":gn('W32'),
    }

# ── Drawing helpers ────────────────────────────────────────────────────────────
def add_rect(slide,x,y,w,h,fill=WHT,lc=None,lw=Pt(0.5)):
    sh=slide.shapes.add_shape(1,I(x),I(y),I(w),I(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else:    sh.fill.background()
    if lc:   sh.line.color.rgb=lc; sh.line.width=lw
    else:    sh.line.fill.background()
    return sh

def add_txt(slide,text,x,y,w,h,sz=9,bold=False,italic=False,
            color=DRK,align=PP_ALIGN.LEFT,face="Calibri"):
    tb=slide.shapes.add_textbox(I(x),I(y),I(w),I(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=str(text)
    r.font.size=P(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name=face
    return tb

def fmt(v,dash_zero=False):
    if v is None: return "-"
    try:
        i=int(v)
        if i==0 and dash_zero: return "-"
        return f"{i:,}" if abs(i)>=1000 else str(i)
    except: return str(v)

def dash(v):
    if not v or v==0: return "-"
    return fmt(v)

def sign(n):
    if n>0: return f"+{fmt(n)}"
    if n<0: return fmt(n)
    return "-"

# ── Cell = rect + textbox (full pixel control, no colspan bugs) ────────────────
def cell(slide,text,x,y,w,h,bg=WHT,color=DRK,sz=7.5,bold=False,
         italic=False,align=PP_ALIGN.CENTER,bc=BDR):
    sh=slide.shapes.add_shape(1,I(x),I(y),I(w),I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=bg
    sh.line.color.rgb=bc; sh.line.width=Pt(0.4)
    tb=slide.shapes.add_textbox(I(x+0.02),I(y+0.01),I(w-0.04),I(h-0.02))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=str(text)
    r.font.size=P(sz); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name="Calibri"

def draw_table(slide, rows, x, y, col_w, row_h, dsz=7.5, bc=BDR):
    """
    rows: list of rows. Each row: list of dicts OR None (skip=spanned cell)
    dict keys: text, bg, color, sz, bold, italic, align, cs, rs
    Spanning: cell with cs=2 spans 2 cols. Mark covered cells as None.
    """
    cy=y
    for ri,row in enumerate(rows):
        cx=x
        rh=row_h[ri] if ri<len(row_h) else row_h[-1]
        for ci,cd in enumerate(row):
            cw=col_w[ci] if ci<len(col_w) else col_w[-1]
            if cd is None:
                cx+=cw; continue
            cs=cd.get('cs',1); rs=cd.get('rs',1)
            tw=sum(col_w[ci:ci+cs])
            th=sum(row_h[ri:ri+rs])
            cell(slide,
                 cd.get('text',''), cx,cy,tw,th,
                 bg=cd.get('bg',WHT), color=cd.get('color',DRK),
                 sz=cd.get('sz',dsz), bold=cd.get('bold',False),
                 italic=cd.get('italic',False),
                 align=cd.get('align',PP_ALIGN.CENTER), bc=bc)
            cx+=cw
        cy+=rh

# ── Cell factories ─────────────────────────────────────────────────────────────
def H(t,cs=1,rs=1,sz=7.5):
    return dict(text=t,bg=NAV,color=WHT,bold=True,sz=sz,cs=cs,rs=rs)
def GH(t,cs=1,rs=1,sz=7.5):
    return dict(text=t,bg=GRY1,color=DRK,bold=False,sz=sz,cs=cs,rs=rs)
def D(t,bg=WHT,color=DRK,bold=False,italic=False,sz=7.5,al=PP_ALIGN.CENTER,cs=1,rs=1):
    return dict(text=str(t) if t not in(None,'') else'-',
                bg=bg,color=color,bold=bold,italic=italic,sz=sz,align=al,cs=cs,rs=rs)
def Y(t,color=DRK,bold=False): return D(t,bg=YLW,color=color,bold=bold)
def B(t,bold=False):            return D(t,color=BLU,bold=bold)
def R(t,bold=False,italic=False): return D(t,color=RED,bold=bold,italic=italic)
def G(t):                       return D(t,color=GRN)
def LB(t):                      return D(t,bg=LBBL,color=BLU)
def E(bg=WHT):                  return D('',bg=bg)

# ── MAIN ───────────────────────────────────────────────────────────────────────
def generate_ppt(file_bytes):
    d=read_excel(file_bytes)
    prs=Presentation()
    prs.slide_width=I(13.33); prs.slide_height=I(7.5)
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide,0,0,13.33,7.5,WHT)

    # ── TITLE ──────────────────────────────────────────────────────────────────
    # Black square icon (hollow)
    add_rect(slide,0.15,0.09,0.22,0.22,BLK)
    add_rect(slide,0.17,0.11,0.18,0.18,WHT)
    add_rect(slide,0.18,0.12,0.15,0.15,BLK)
    # Title text
    add_txt(slide,"Engine Daily Production Report _ Assy",
            0.46,0.06,8.5,0.36,sz=24,bold=True,color=BLK,align=PP_ALIGN.LEFT)
    # Date
    add_txt(slide,d['date_str'],
            9.2,0.06,4.0,0.36,sz=14,color=BLU,align=PP_ALIGN.RIGHT)
    # Underline
    add_rect(slide,0.10,0.44,13.13,0.012,BLK)

    LX=0.10; LW=6.22
    RX=6.65; RW=6.58

    # ─── EFFICIENCY TRENDS ─────────────────────────────────────────────────────
    add_rect(slide,LX,0.50,0.14,0.14,BLK)
    add_txt(slide," Efficiency Trends",LX+0.16,0.47,3.5,0.20,
            sz=11,bold=True,color=BLK)

    # Chart outer box
    BX=LX; BY=0.68; BW=LW; BH=2.80
    add_rect(slide,BX,BY,BW,BH,WHT,BLK,Pt(0.75))

    # KPI panel left of chart
    KX=BX+0.06
    add_txt(slide,"□ Overall",KX,BY+0.05,1.55,0.17,sz=7.5,color=DRK)
    add_txt(slide,f"  : {d['eff_c']}",KX,BY+0.20,1.55,0.17,sz=7.5,color=DRK)
    add_txt(slide,f"✱ UPH: {d['uph']}",KX,BY+0.40,1.55,0.20,sz=9.5,bold=True,color=DRK)
    add_txt(slide,"✱ Work'g Time",KX,BY+0.62,1.55,0.20,sz=9.5,bold=True,color=DRK)
    add_txt(slide,f"  \u2212 1Sh : {dash(d['pt_a'])}'",KX,BY+0.84,1.55,0.18,sz=8.5,color=DRK)
    add_txt(slide,f"  \u2212 2Sh : {dash(d['pt_b'])}'",KX,BY+1.02,1.55,0.18,sz=8.5,color=DRK)
    add_txt(slide,f"  \u2212 3Sh : {dash(d['pt_c'])}'",KX,BY+1.20,1.55,0.18,sz=8.5,color=DRK)

    # Real bar chart
    date_obj=d['date_obj']
    if date_obj:
        from datetime import timedelta
        mon=date_obj-timedelta(days=date_obj.weekday())
        dnames=["(Mon)","(Tue)","(Wed)","(Thu)","(Fri)","(Sat)","(Sun)"]
        daily_labels=[f"{(mon+timedelta(days=i)).day:02d}/{(mon+timedelta(days=i)).strftime('%m')}\n{dnames[i]}" for i in range(7)]
        today_col=date_obj.weekday()
    else:
        daily_labels=["23/02\n(Mon)","24/02\n(Tue)","25/02\n(Wed)",
                      "26/02\n(Thu)","27/02\n(Fri)","28/02\n(Sat)","01/03\n(Sun)"]
        today_col=1

    cats=["'26y\nCumu.","Jan.26","Feb.26\nCumu.","Last\nWeek","This\nWeek"]+daily_labels
    eff_val=d['eff_c_val'] if d['eff_c_val'] else 100.0
    vals=[100.0,100.0,100.0,100.0,100.0]+[eff_val if i==today_col else (None if i>today_col else 100.0) for i in range(7)]

    cd2=ChartData(); cd2.categories=cats; cd2.add_series('Efficiency',vals)
    cshape=slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        I(BX+1.60),I(BY+0.04),I(BW-1.65),I(BH-0.10),cd2)
    ch=cshape.chart
    ch.has_legend=False; ch.has_title=False
    pass  # chart area

    va=ch.value_axis
    va.minimum_scale=0; va.maximum_scale=120; va.major_unit=20
    va.has_major_gridlines=True
    va.major_gridlines.format.line.color.rgb=RGBColor(0xD0,0xD0,0xD0)
    va.tick_labels.font.size=P(7); va.tick_labels.font.color.rgb=DRK
    ca=ch.category_axis
    ca.has_major_gridlines=False
    ca.tick_labels.font.size=P(6); ca.tick_labels.font.color.rgb=DRK

    ser=ch.series[0]
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb=RGBColor(0x1F,0x50,0x7D)
    ser.data_labels.show_value=True
    ser.data_labels.font.size=P(7); ser.data_labels.font.bold=True
    ser.data_labels.font.color.rgb=BLU

    # Red dashed reference line at 100 (drawn as thin red rect)
    chart_inner_y=BY+0.41
    add_rect(slide,BX+1.66,chart_inner_y,BW-1.72,0.014,RED)

    # Orange border around today's bar
    n_cats=len(cats)
    bar_zone_w=BW-1.72
    bar_w=bar_zone_w/n_cats
    today_idx=5+today_col
    box_x=BX+1.66+today_idx*bar_w-0.01
    add_rect(slide,box_x,BY+0.06,bar_w+0.02,BH-0.14,None,RGBColor(0xFF,0x80,0x00),Pt(1.5))

    # ─── EFFICIENCY TABLE ──────────────────────────────────────────────────────
    ETY=3.54
    add_rect(slide,LX,ETY,0.14,0.14,BLK)
    add_txt(slide," Efficiency",LX+0.16,ETY-0.03,2.5,0.20,sz=11,bold=True,color=BLK)

    # 13 cols
    C=[0.56,0.60,0.46,0.44,0.46,0.44,0.44,0.44,0.42,0.42,0.42,0.42,0.42]
    C[12]=LW-sum(C[:12])
    RH=0.215; RHH=0.24

    this_in=(d['k10t_in_m']+d['ak10m_in_m']+d['ak12m_in_m']+
             d['gege_in_m']+d['geng_in_m'])
    this_out=(d['k10t_out_m']+d['ak10m_out_m']+d['ak12m_out_m']+
              d['gege_out_m']+d['geng_out_m'])
    today_in=d['ttl_in_c']; today_out=d['ttl_out_c']
    stk=d['ttl_stk']

    C26=56782;Cjan=31292;Cfeb=d['ttl_feb_p'];Clast=7794
    R26=56824;Rjan=31315;Rfeb=d['ttl_feb_r'];Rlast=7800

    # day column headers
    if date_obj:
        from datetime import timedelta
        mon2=date_obj-timedelta(days=date_obj.weekday())
        dn=["(Mon)","(Tue)","(Wed)","(Thu)","(Fri)","(Sat)","(Sun)"]
        dh=[f"{(mon2+timedelta(days=i)).day:02d}/{(mon2+timedelta(days=i)).strftime('%m')}\n{dn[i]}" for i in range(7)]
        tc=date_obj.weekday()
    else:
        dh=["23/02\n(Mon)","24/02\n(Tue)","25/02\n(Wed)",
            "26/02\n(Thu)","27/02\n(Fri)","28/02\n(Sat)","01/03\n(Sun)"]
        tc=1

    # Build 5 day columns: today + next 4
    day_hdrs=[dh[(tc+i)%7] for i in range(5)]

    def day_v(val, i):
        if i==0: return Y(fmt(val),BLU)     # today=yellow+blue
        else:    return E(GRY1)              # future=gray empty

    tbl=[
        # Header
        [GH("Category\n4,025",cs=2,sz=7),None,
         GH("'26y\nCumu.",sz=7),GH("Jan.26",sz=7),GH("Feb.26\nCumu.",sz=7),GH("Last\nWeek",sz=7),
         GH("This\nWeek",sz=7),GH(day_hdrs[0],sz=6.5),
         GH(day_hdrs[1],sz=6.5),GH(day_hdrs[2],sz=6.5),
         GH(day_hdrs[3],sz=6.5),GH(day_hdrs[4],sz=6.5),
         GH("",sz=6)],
        # Capa
        [GH("Prod\n(Q'ty)",rs=5),D("Capa.",bg=GRY1),
         D(fmt(C26)),D(fmt(Cjan)),D(fmt(Cfeb)),D(fmt(Clast)),
         Y(fmt(this_in)),*[day_v(today_in,i) for i in range(5)],E(GRY1)],
        # Plan
        [None,D("Plan",bg=GRY1),
         D(fmt(C26)),D(fmt(Cjan)),D(fmt(Cfeb)),D(fmt(Clast)),
         Y(fmt(this_in)),*[day_v(today_in,i) for i in range(5)],E(GRY1)],
        # Result
        [None,D("Result",bg=LBBL,color=BLU,bold=True),
         LB(fmt(R26)),LB(fmt(Rjan)),LB(fmt(Rfeb)),LB(fmt(Rlast)),
         Y(fmt(this_out),BLU),*[day_v(today_out,i) for i in range(5)],E(GRY1)],
        # Diff
        [None,D("Diff.",bg=GRY1),
         G(sign(R26-C26)),G(sign(Rjan-Cjan)),G(sign(Rfeb-Cfeb)),G(sign(Rlast-Clast)),
         Y(sign(this_out-this_in),GRN),Y("-"),
         E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        # Stock
        [None,Y("Stock",DRK,True),
         D(fmt(4246)),D(fmt(6268)),D(fmt(4246)),D(fmt(4448)),
         Y(fmt(this_out)),*[day_v(stk,i) for i in range(5)],E(GRY1)],
        # Work
        [GH("Time\n(Min)",rs=3),D("Work",bg=GRY1),
         D(fmt(61954)),D(fmt(34142)),D(fmt(d['apt_month'])),D(fmt(8504)),
         Y(fmt(d['apt_month'])),*[day_v(d['apt_c'],i) for i in range(5)],E(GRY1)],
        [None,D("Plan.Stop",bg=GRY1),
         D('-'),D('-'),D('-'),D('-'),Y('-'),Y('-'),
         E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        [None,D("Operation",bg=GRY1),
         D(fmt(61954)),D(fmt(34142)),D(fmt(d['apt_month'])),D(fmt(8504)),
         Y(fmt(d['apt_month'])),*[day_v(d['apt_c'],i) for i in range(5)],E(GRY1)],
        # Down time
        [GH("Down\nTime\n(Min)",rs=6),D("Equip.",bg=GRY1),
         D(dash(d['dt_equip'])),D('-'),D('-'),D('-'),
         Y(dash(d['dt_equip'])),*[day_v(d['dt_equip'],i) for i in range(5)],E(GRY1)],
        [None,D("No parts",bg=GRY1),
         D(dash(d['dt_nopart'])),D('-'),D('-'),D('-'),
         Y(dash(d['dt_nopart'])),Y('-'),E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        [None,D("NG parts",bg=GRY1),
         D('-'),D('-'),D('-'),D('-'),Y('-'),Y('-'),
         E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        [None,D("Work Delay",bg=GRY1),
         D('-'),D('-'),D('-'),D('-'),Y('-'),Y('-'),
         E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        [None,D("M/C Sett'g",bg=GRY1),
         D('-'),D('-'),D('-'),D('-'),Y('-'),Y('-'),
         E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        [None,D("Others",bg=GRY1),
         D('-'),D('-'),D('-'),D('-'),Y('-'),Y('-'),
         E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        # Sub total
        [E(GRY1),D("Sub Total",bg=GRY1),
         D(dash(d['dt_ttl'])),D('-'),D('-'),D('-'),
         Y(dash(d['dt_ttl'])),Y(dash(d['dt_ttl']),BLU),
         E(GRY1),E(GRY1),E(GRY1),E(GRY1),E(GRY1)],
        # Overall Effi
        [Y("Overall Effi.",DRK,True),None,
         Y(d['eff_m'],BLU,True),Y(d['eff_m'],BLU,True),
         Y(d['eff_m'],BLU,True),Y(d['eff_m'],BLU,True),
         Y(d['eff_c'],BLU,True),Y(d['eff_c'],BLU,True),
         Y(''),Y(''),Y(''),Y(''),Y('')],
    ]

    rhl=[RHH]+[RH]*(len(tbl)-1)
    draw_table(slide,tbl,LX,ETY+0.17,C,rhl,bc=BDR)

    # ─── PROD & STOCK STATUS ───────────────────────────────────────────────────
    add_rect(slide,RX,0.50,0.14,0.14,BLK)
    add_txt(slide," Prod. & Stock Status",RX+0.16,0.47,4.5,0.20,sz=11,bold=True,color=BLK)

    SC=[0.50,0.38,0.38,0.38,0.52,0.52,0.44,0.50,0.44,0.44,0.50,0.55]
    SC[11]=RW-sum(SC[:11])
    SRH=0.20

    stbl=[
        [H("Model"),H("Prod.",cs=3),None,None,
         H("Feb."),H(""),H("Total\n('26)"),
         H("Stock",cs=4),None,None,None],
        [E(NAV),H("3 Sh"),H("1 Sh"),H("2 Sh"),
         H("Plan"),H("Result"),H("%"),
         H("Engine",sz=6.5),H("Mobis",sz=6.5),
         H("Glovis",sz=6.5),H("Ext.WH",sz=6.5),H("Sum")],
        # K1.0T
        [D("K1.0T"),B(fmt(d['k10t_in_c'])),D("-"),D("-"),
         D(fmt(d['k10t_feb_p'],True)),D(fmt(d['k10t_feb_r'],True)),
         R(d['k10t_pct'],bold=True,italic=True),
         D(fmt(1590,True)),D(fmt(170,True)),D(fmt(32,True)),D(fmt(522,True)),
         B(fmt(d['k10t_stk']))],
        # K1.0M
        [D("K1.0M"),D("-"),D("-"),D("-"),
         D("-"),D("-"),D("-"),
         D("-"),D("-"),D("-"),D("-"),D("-")],
        # AK 1.2 MPI
        [D("AK\n1.2",rs=2),D("MPI"),B(fmt(d['ak12m_in_c'])),D("-"),
         D(fmt(d['ak12m_feb_p'],True)),D(fmt(d['ak12m_feb_r'],True)),
         R(d['ak12m_pct'],bold=True,italic=True),
         D(fmt(13039,True)),D(fmt(389,True)),D(fmt(50,True)),D(fmt(606,True)),
         B(fmt(d['ak12m_stk']))],
        [None,D("CVT"),D("-"),D("-"),D("-"),D("-"),D("-"),
         D("-"),D("-"),D("-"),D("-"),D("-")],
        # GII EGR
        [D("GII",rs=2),D("EGR"),B(fmt(d['gege_in_c'])),D("-"),
         D(fmt(d['gege_feb_p'],True)),D(fmt(d['gege_feb_r'],True)),
         R(d['gege_pct'],bold=True,italic=True),
         D(fmt(30708,True)),D(fmt(678,True)),D(fmt(143,True)),D(fmt(1236,True)),
         B(fmt(d['gege_stk']))],
        [None,D("-EGR"),B(fmt(d['geng_in_c'])),D("-"),
         D(fmt(d['geng_feb_p'],True)),D(fmt(d['geng_feb_r'],True)),
         R(d['geng_pct'],bold=True,italic=True),
         D(fmt(11487,True)),D(fmt(186,True)),D(fmt(66,True)),D(fmt(168,True)),
         B(fmt(d['geng_stk']))],
        # Total
        [D("Total",bold=True,cs=2),None,
         D(fmt(d['ttl_in_c'])),
         D(fmt(d['ttl_feb_p'],True)),D(fmt(d['ttl_feb_r'],True)),
         R(d['ttl_pct'],bold=True,italic=True),
         D(fmt(R26,True)),
         D(fmt(1423,True)),D(fmt(291,True)),D(fmt(2532,True)),
         D("-"),B(fmt(d['ttl_stk']))],
    ]
    srhl=[SRH*0.85,SRH*0.85]+[SRH]*7
    draw_table(slide,stbl,RX,0.68,SC,srhl,bc=BDR)

    # ─── DOWN TIME ─────────────────────────────────────────────────────────────
    DTY=3.54
    add_rect(slide,RX,DTY,0.14,0.14,BLK)
    add_txt(slide," Down Time",RX+0.16,DTY-0.03,3.0,0.20,sz=11,bold=True,color=BLK)

    DC=[0.88,3.40,0.85,0.85]
    DC[1]=RW-DC[0]-DC[2]-DC[3]
    DTRH=0.285; DTRHH=0.26

    dtv=lambda v: str(v) if v and v!=0 else "-"

    dt_items=[
        ("Equip\nTrouble",d['dt_equip']),
        ("No\nParts",     d['dt_nopart']),
        ("NG\nParts",     d['dt_ngpart']),
        ("Works\nDelay",  d['dt_wdlay']),
        ("M/C\nSetting",  d['dt_mcset']),
        ("Others",        d['dt_other']),
    ]
    dtbl=[[H("Items",sz=9),H("Major issues",sz=9),H("Resp.",sz=9),H("D/T\n(Min)",sz=9)]]
    for item,val in dt_items:
        dtbl.append([D(item),D("-"),D("-"),D(dtv(val))])
    dtbl.append([Y("Total",DRK,True),Y("-"),Y("-"),Y(dtv(d['dt_ttl']))])
    dtbl.append([D("Planned\nStop"),D("-"),D("-"),D("-")])

    nd=len(dtbl)
    rem=7.42-(DTY+0.17)
    dtrhl=[DTRHH]+[(rem-DTRHH)/(nd-1)]*(nd-1)
    draw_table(slide,dtbl,RX,DTY+0.17,DC,dtrhl,dsz=8.5,bc=BDR)

    out=io.BytesIO(); prs.save(out); out.seek(0)
    return out.read()